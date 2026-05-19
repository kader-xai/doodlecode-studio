"""Mini-tools that live outside the slide deck — used by /tools.

First tool: PPT → Images.
  - Uploads a .pptx
  - Renders each slide to a PNG (LibreOffice headless if available;
    falls back to a python-pptx text-only renderer)
  - Extracts speaker notes into a notes.txt
  - Output folder lives under ~/.doodlecode/tools/<deck_name>/
"""
from __future__ import annotations

import re
import shutil
import subprocess
import tempfile
import textwrap
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Optional


TOOLS_DIR = Path.home() / ".doodlecode" / "tools"
TOOLS_DIR.mkdir(parents=True, exist_ok=True)


@dataclass
class PptResult:
    deck_name: str
    folder: str          # absolute path on disk
    slides: list[str]    # filenames inside folder (slide_001.png, ...)
    notes_file: str      # notes.txt filename
    renderer: str        # "libreoffice" or "python-pptx-fallback"
    message: str = ""    # info / hint


def _slugify(name: str) -> str:
    """Filesystem-safe folder name derived from the deck filename."""
    base = Path(name).stem or "deck"
    s = re.sub(r"[^A-Za-z0-9._-]+", "_", base).strip("_")
    return s or "deck"


def _find_libreoffice() -> Optional[str]:
    for cand in ("soffice", "libreoffice"):
        p = shutil.which(cand)
        if p:
            return p
    return None


def _libreoffice_render(pptx_path: Path, out_dir: Path) -> list[str]:
    """LibreOffice → PDF, then pdftoppm → PNGs. Returns slide filenames."""
    soffice = _find_libreoffice()
    if not soffice:
        raise FileNotFoundError("libreoffice / soffice not found on PATH")
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        raise FileNotFoundError("pdftoppm (poppler) not found on PATH")
    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        # PPTX → PDF (in a sandbox dir so user-profile noise stays put)
        subprocess.run(
            [soffice, "--headless", "--norestore", "--nologo", "--nodefault",
             "--convert-to", "pdf", "--outdir", str(td_path), str(pptx_path)],
            check=True, capture_output=True, timeout=180,
        )
        pdfs = list(td_path.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice produced no PDF")
        pdf_path = pdfs[0]
        # PDF → PNGs (one per page). pdftoppm prefix-{page}.png
        prefix = out_dir / "slide"
        subprocess.run(
            [pdftoppm, "-png", "-r", "150", str(pdf_path), str(prefix)],
            check=True, capture_output=True, timeout=180,
        )
    # pdftoppm names files like slide-1.png, slide-2.png … rename to
    # slide_001.png so they sort lexically forever.
    renamed: list[str] = []
    raw = sorted(out_dir.glob("slide-*.png"),
                 key=lambda p: int(re.search(r"-(\d+)\.png$", p.name).group(1)))
    for i, p in enumerate(raw, start=1):
        target = out_dir / f"slide_{i:03d}.png"
        p.rename(target)
        renamed.append(target.name)
    return renamed


def _pptx_fallback_render(pptx_path: Path, out_dir: Path) -> list[str]:
    """No LibreOffice — render each slide as a text-only Pillow image
    listing the slide title + text frames + a 'Speaker notes' footer.
    Not a perfect visual reproduction, but useful for note-taking."""
    from pptx import Presentation  # type: ignore
    from PIL import Image, ImageDraw, ImageFont  # type: ignore

    prs = Presentation(str(pptx_path))
    out_files: list[str] = []
    W, H = 1600, 900
    bg, fg = "white", "#1a1a1a"
    try:
        title_font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial Bold.ttf", 48)
        body_font  = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 28)
    except OSError:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    for i, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGB", (W, H), bg)
        d = ImageDraw.Draw(img)
        y = 60
        # Title from first non-empty text frame, then body bullets.
        body_lines: list[str] = []
        title_text = f"Slide {i}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs).strip()
                if not txt:
                    continue
                if title_text == f"Slide {i}":
                    title_text = txt
                else:
                    body_lines.append(txt)
        d.text((80, y), title_text[:80], font=title_font, fill=fg)
        y += 90
        for line in body_lines:
            for wrap in textwrap.wrap(line, width=70):
                d.text((100, y), "• " + wrap, font=body_font, fill=fg)
                y += 40
                if y > H - 80:
                    break
            if y > H - 80:
                break
        # Footer
        d.text((80, H - 50), f"Slide {i}  ·  fallback render (LibreOffice not installed)",
               font=body_font, fill="#888")
        name = f"slide_{i:03d}.png"
        img.save(out_dir / name)
        out_files.append(name)
    return out_files


def _extract_notes(pptx_path: Path) -> list[tuple[int, str, str]]:
    """Returns [(slide_index, title, notes_text)]."""
    from pptx import Presentation  # type: ignore
    prs = Presentation(str(pptx_path))
    out: list[tuple[int, str, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    title = t.split("\n", 1)[0][:120]
                    break
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        out.append((i, title, notes))
    return out


def _write_notes_file(folder: Path, notes: list[tuple[int, str, str]]) -> str:
    name = "notes.txt"
    lines: list[str] = []
    lines.append(f"# Speaker notes — {folder.name}")
    lines.append(f"# {len(notes)} slides")
    lines.append("")
    for i, title, body in notes:
        lines.append(f"## Slide {i}: {title or '(untitled)'}")
        if body:
            lines.append(body)
        else:
            lines.append("(no notes)")
        lines.append("")
    (folder / name).write_text("\n".join(lines), encoding="utf-8")
    return name


def convert_pptx(filename: str, data: bytes) -> PptResult:
    if not filename.lower().endswith((".pptx", ".pptm")):
        raise ValueError("Only .pptx / .pptm files are supported.")
    # Sanity-check: pptx is a zip with ppt/ inside.
    try:
        zipfile.ZipFile(__import__("io").BytesIO(data)).namelist()
    except zipfile.BadZipFile as e:
        raise ValueError("File is not a valid .pptx (not a zip archive).") from e

    deck_name = _slugify(filename)
    folder = TOOLS_DIR / deck_name
    if folder.exists():
        shutil.rmtree(folder)
    folder.mkdir(parents=True, exist_ok=True)

    src = folder / filename
    src.write_bytes(data)

    renderer = "libreoffice"
    message = ""
    try:
        slides = _libreoffice_render(src, folder)
    except FileNotFoundError as e:
        renderer = "python-pptx-fallback"
        message = (
            f"{e}. Used the python-pptx fallback renderer (text-only). "
            "Install LibreOffice for accurate visuals: "
            "`brew install --cask libreoffice`."
        )
        slides = _pptx_fallback_render(src, folder)

    notes_pairs = _extract_notes(src)
    notes_file = _write_notes_file(folder, notes_pairs)

    return PptResult(
        deck_name=deck_name,
        folder=str(folder),
        slides=slides,
        notes_file=notes_file,
        renderer=renderer,
        message=message,
    )
